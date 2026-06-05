#!/usr/bin/env python3
"""pptx-extract.py — extract slide text from a .pptx clue into a flat text file
mirroring `pdftotext -layout` output (slide-ordered; text frames + tables).
Used when a consulting-toolkit clue arrives as PowerPoint instead of PDF.

Usage: DYLD_FALLBACK_LIBRARY_PATH=/opt/homebrew/lib python3 pptx-extract.py <in.pptx> <out.txt>
Requires: python-pptx.
"""
import sys
from pptx import Presentation

def text_of(shape, out):
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            t = "".join(r.text for r in para.runs)
            if t.strip():
                out.append(t)
    if shape.shape_type == 6:  # group
        for sub in shape.shapes:
            text_of(sub, out)
    if getattr(shape, "has_table", False) and shape.has_table:
        for row in shape.table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                out.append(" | ".join(cells))

def main():
    src, dst = sys.argv[1], sys.argv[2]
    prs = Presentation(src)
    out = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"\n===== SLIDE {i} =====")
        for shape in slide.shapes:
            text_of(shape, out)
    with open(dst, "w") as f:
        f.write("\n".join(out) + "\n")
    words = len(" ".join(out).split())
    print(f"extracted {len(prs.slides._sldIdLst)} slides, ~{words} words -> {dst}")

if __name__ == "__main__":
    main()
