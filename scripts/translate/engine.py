#!/usr/bin/env python3
"""engine.py — ingest + structure + split + assemble + plan + brief for /translate.
Subcommands:
  plan     <workdir>   read workdir/config.json -> ingest src, detect mode, split into
                       translatable units, write brief.md + plan.json + meta.json
  assemble <workdir>   read translated units_<lang>/ -> book_blocks.json (for render.py)
Multi-format ingest: pdf, docx, pptx, html/url, md, txt. Structure-aware (book / doc /
slides). Books auto-split into chapters; long chapters split at section boundaries.
"""
import sys, os, re, json, html as _html, collections, math
from pathlib import Path

CAP = 3500  # max source words per translation unit (keeps each parallel agent world-class)

# ----------------------------------------------------------------------------- utils
def slugify(s, n=48):
    s = re.sub(r"[^\w\s-]", "", (s or "").lower(), flags=re.U)
    s = re.sub(r"[\s_-]+", "-", s).strip("-")
    return (s[:n] or "untitled")

def wc(s): return len(re.findall(r"\w+", s, re.U))

def dehyph(lines):
    out = ""
    for t in lines:
        if out.endswith("-") and re.search(r"[^\W\d_]-$", out, re.U) and re.match(r"^[a-zà-ỹ]", t, re.I):
            out = out[:-1] + t.lstrip()
        else:
            out = (out + " " + t).strip() if out else t
    return re.sub(r"[ \t]+", " ", out).strip()

# ----------------------------------------------------------------------------- PDF
# v0.3 STEM figure/table fidelity (capability `translate` extend). PyMuPDF's
# page.get_images() only returns EMBEDDED RASTER xrefs — so VECTOR figures (matplotlib
# violin plots, learning curves, parameter studies) are invisible to it and were dropped,
# leaving caption-only gaps. Instead we anchor on each "Figure/Table/Hình/Bảng N" caption
# and render the adjacent page REGION to a PNG via get_pixmap(clip=rect), capturing vector
# AND raster faithfully. --math=crop reuses the same machinery for display equations.
def _area(r):
    return max(0.0, r.x1 - r.x0) * max(0.0, r.y1 - r.y0)

def _inter_area(a, b):
    return (max(0.0, min(a.x1, b.x1) - max(a.x0, b.x0))
            * max(0.0, min(a.y1, b.y1) - max(a.y0, b.y0)))

def _overlap_frac(a, b):
    """Fraction of rect a that lies inside rect b."""
    aa = _area(a)
    return (_inter_area(a, b) / aa) if aa > 0 else 0.0

def _union(rects):
    import fitz
    return fitz.Rect(min(r.x0 for r in rects), min(r.y0 for r in rects),
                     max(r.x1 for r in rects), max(r.y1 for r in rects))

def _graphic_rects(page):
    """Bounding rects of vector drawings + raster images — the figure-bearing elements.
    Vector drawings are exactly what page.get_images() misses."""
    import fitz
    out = []
    try:
        for d in page.get_drawings():
            r = fitz.Rect(d["rect"])
            if r.width > 3 and r.height > 3:
                out.append(r)
    except Exception:
        pass
    try:
        for im in page.get_images(full=True):
            for r in page.get_image_rects(im[0]):
                out.append(fitz.Rect(r))
    except Exception:
        pass
    return out

_CAPTION_RE = re.compile(r"^\s*(figure|fig\.?|hình|table|bảng)\s*\d", re.I)

def _caption_figures(page, content):
    """Caption-anchored figure/table regions: [{rect, cap_rect, cap_y, caption, kind}].
    The crop is bounded by the nearest WIDE body paragraph (above the caption for figures,
    below it for tables) so it spans the whole graphic (incl. stacked subplots) yet can
    never swallow prose; it spans the content width so axis labels are both captured in the
    image and suppressed from the text stream."""
    import fitz
    pr = page.rect
    blocks = [b for b in page.get_text("dict")["blocks"] if b.get("lines")]
    def _chars(b):
        return sum(len(s["text"]) for l in b["lines"] for s in l["spans"])
    # body paragraphs = wide + substantial text (figure axis labels are short → excluded)
    body_blocks = [fitz.Rect(b["bbox"]) for b in blocks
                   if (b["bbox"][2] - b["bbox"][0]) > content.width * 0.5 and _chars(b) > 80]
    caps = []
    for b in blocks:
        first = "".join(s["text"] for s in b["lines"][0]["spans"]).strip()
        if _CAPTION_RE.match(first):
            full = " ".join("".join(s["text"] for s in l["spans"]) for l in b["lines"])
            kind = "table" if re.match(r"^\s*(table|bảng)", first, re.I) else "figure"
            caps.append((fitz.Rect(b["bbox"]), re.sub(r"\s+", " ", full).strip(), kind))
    if not caps:
        return []
    graphics = _graphic_rects(page)
    out = []
    for cap_rect, cap_text, kind in caps:
        # band = the vertical gap between the caption and the nearest body paragraph on the
        # figure's side. Tables sit below their caption; figures above.
        if kind == "table":
            bound = min([r.y0 for r in body_blocks if r.y0 >= cap_rect.y1 - 2] or [pr.y1])
            lo, hi = cap_rect.y1, bound
        else:
            bound = max([r.y1 for r in body_blocks if r.y1 <= cap_rect.y0 + 2] or [pr.y0])
            lo, hi = bound, cap_rect.y0
        # every graphic overlapping the band, clipped to it (no width filter — vector plots
        # render as many thin path segments). The union's y-extent is the figure's height.
        clipped = [(max(g.y0, lo), min(g.y1, hi)) for g in graphics if g.y1 > lo + 1 and g.y0 < hi - 1]
        clipped = [(a, b) for a, b in clipped if b - a > 1]
        if not clipped:
            continue
        y0 = max(pr.y0, min(a for a, _ in clipped))
        y1 = min(pr.y1, max(b for _, b in clipped))
        if (y1 - y0) < pr.height * 0.045:
            continue
        reg = fitz.Rect(content.x0, y0, content.x1, y1)
        out.append({"rect": reg, "cap_rect": cap_rect, "cap_y": cap_rect.y0,
                    "caption": cap_text.replace("]", ")").strip(), "kind": kind})
    return out

_MATHY_RE = re.compile(r"[=≈≤≥<>∑∏∫√±×·∂πΑ-Ωα-ω⌘⇤⇣→]|\\frac|\^|_\{|\.=")

def _equation_crops(page, body_size, content):
    """--math=crop only: display-equation regions (short, horizontally inset/centered,
    math-symbol-dense or carrying a trailing (N.M) number). Best-effort heuristic."""
    import fitz
    out = []
    for b in page.get_text("dict")["blocks"]:
        if not b.get("lines") or len(b["lines"]) > 6:
            continue
        bb = fitz.Rect(b["bbox"])
        txt = " ".join("".join(s["text"] for s in l["spans"]) for l in b["lines"]).strip()
        if not (2 <= len(txt) <= 260):
            continue
        inset = bb.x0 > content.x0 + content.width * 0.08 and bb.x1 < content.x1 - content.width * 0.02
        has_num = bool(re.search(r"\(\d+\.\d+[a-z]?\)\s*$", txt))
        sym = sum(1 for c in txt if not (c.isalnum() or c.isspace()))
        dense = sym / max(1, len(txt)) > 0.16
        if (has_num or (inset and dense)) and _MATHY_RE.search(txt):
            m = re.search(r"\((\d+\.\d+[a-z]?)\)\s*$", txt)
            out.append({"rect": bb, "num": m.group(1) if m else None})
    return out

def ingest_pdf(path, assets_dir=None, keep_assets=True, math_mode="auto"):
    import fitz
    doc = fitz.open(path)
    meta = {"title": (doc.metadata or {}).get("title") or "", "author": (doc.metadata or {}).get("author") or "",
            "format": "pdf", "n": doc.page_count}
    toc = doc.get_toc() or []
    if keep_assets and assets_dir:
        Path(assets_dir).mkdir(parents=True, exist_ok=True)
    # body font size = modal size weighted by char count
    sizes = collections.Counter()
    for p in doc:
        for b in p.get_text("dict")["blocks"]:
            for l in b.get("lines", []):
                for s in l["spans"]:
                    if s["text"].strip():
                        sizes[round(s["size"], 1)] += len(s["text"].strip())
    body = sizes.most_common(1)[0][0] if sizes else 10.0
    # repeating header/footer lines (same short text near top/bottom on many pages)
    edge = collections.Counter()
    for p in doc:
        h = p.rect.height
        for b in p.get_text("dict")["blocks"]:
            for l in b.get("lines", []):
                y = l["bbox"][1]
                txt = "".join(s["text"] for s in l["spans"]).strip()
                if txt and (y < h * 0.09 or y > h * 0.91) and len(txt) < 70:
                    edge[re.sub(r"\d+", "#", txt)] += 1
    strip = {k for k, v in edge.items() if v >= max(3, doc.page_count * 0.25)}

    # chapter start pages from TOC level-1
    toc_starts = {e[2] - 1: e[1] for e in toc if e[0] <= 1} if toc else {}

    out = []
    nfig = 0
    zoom_fig = fitz.Matrix(2.8, 2.8)   # ~200 dpi region render (vector + raster faithfully)
    zoom_eq = fitz.Matrix(3.4, 3.4)    # crisper for --math=crop equation crops
    for pno, p in enumerate(doc):
        if pno in toc_starts:
            out.append(f"\n# {toc_starts[pno].strip()}\n")
        pr = p.rect
        cmargin = pr.width * 0.04
        content = fitz.Rect(pr.x0 + cmargin, pr.y0, pr.x1 - cmargin, pr.y1)
        items = []        # (y_top, markdown) — text, tables, figures interleaved by position
        crop_rects = []   # regions whose leaked text must be suppressed (now in a crop image)

        # caption-anchored figure/table crops (Req 1): render the page REGION so vector
        # plots are captured, not lost like the embedded-raster path.
        if keep_assets and assets_dir:
            for fr in _caption_figures(p, content):
                nfig += 1
                fn = f"fig-{pno + 1:03d}-{nfig:02d}.png"
                try:
                    p.get_pixmap(matrix=zoom_fig, clip=fr["rect"]).save(str(Path(assets_dir) / fn))
                except Exception:
                    continue
                items.append((fr["cap_y"], f"\n![{fr['caption']}](assets/{fn})\n"))
                crop_rects.append(fr["rect"])
                crop_rects.append(fr["cap_rect"])   # caption now lives in the alt text
            if math_mode == "crop":
                for er in _equation_crops(p, body, content):
                    nfig += 1
                    fn = f"eq-{pno + 1:03d}-{nfig:02d}.png"
                    try:
                        p.get_pixmap(matrix=zoom_eq, clip=er["rect"]).save(str(Path(assets_dir) / fn))
                    except Exception:
                        continue
                    alt = f"Eq. {er['num']}" if er.get("num") else ""
                    items.append((er["rect"].y0, f"\n![{alt}](assets/{fn})\n"))
                    crop_rects.append(er["rect"])

        # grid-lined tables via find_tables, unless they overlap a cropped figure/equation
        table_bboxes = []
        if keep_assets:
            try:
                for tb in p.find_tables().tables:
                    tbr = fitz.Rect(tb.bbox)
                    if any(_overlap_frac(tbr, cr) > 0.5 for cr in crop_rects):
                        continue
                    mt = tb.to_markdown()
                    if mt and mt.count("|") >= 4:
                        items.append((tb.bbox[1], "\n" + mt.strip() + "\n"))
                        table_bboxes.append(tb.bbox)
            except Exception:
                pass

        blocks = [b for b in p.get_text("dict")["blocks"] if "lines" in b]
        for b in blocks:
            bb = fitz.Rect(b["bbox"])
            by0, by1 = b["bbox"][1], b["bbox"][3]
            if any(_overlap_frac(bb, cr) > 0.55 for cr in crop_rects):
                continue  # leaked figure/equation text — captured by the crop image
            if any(tb[1] - 2 <= by0 and by1 <= tb[3] + 2 for tb in table_bboxes):
                continue  # text already captured by the table markdown
            infos = []
            for l in b["lines"]:
                spans = l["spans"]
                real = [s for s in spans if s["text"].strip()]
                if not real:
                    continue
                txt = ""
                for s in spans:
                    tx = s["text"]
                    if (s["flags"] & 1) and tx.strip().isdigit():
                        txt += f"<sup>{tx.strip()}</sup>"
                    else:
                        txt += tx
                dom = max(real, key=lambda s: len(s["text"]))
                infos.append({"t": txt, "raw": "".join(s["text"] for s in spans).strip(),
                              "sz": dom["size"], "bold": bool(dom["flags"] & 16)})
            infos = [i for i in infos if i["raw"] and re.sub(r"\d+", "#", i["raw"]) not in strip
                     and not re.fullmatch(r"\d{1,4}", i["raw"])]
            if not infos:
                continue
            mx = max(i["sz"] for i in infos)
            bold = all(i["bold"] for i in infos)
            text = dehyph([i["t"] for i in infos])
            short = len(" ".join(i["raw"] for i in infos)) < 70
            if mx >= body * 1.7 and short and pno not in toc_starts:
                items.append((by0, f"\n# {text}\n"))
            elif mx >= body * 1.28 and short:
                items.append((by0, f"\n## {text}\n"))
            elif bold and short and mx >= body * 0.95:
                items.append((by0, f"\n## {text}\n"))
            else:
                items.append((by0, text + "\n"))
        items.sort(key=lambda t: t[0])
        for _, m in items:
            out.append(m)
    md = re.sub(r"\n{3,}", "\n\n", "\n".join(out)).strip() + "\n"
    n_h1 = len(re.findall(r"(?m)^# ", md))
    meta["is_book_hint"] = bool(len(toc) >= 3 or doc.page_count >= 40 or n_h1 >= 4)
    if not meta["title"]:
        m = re.search(r"(?m)^# (.+)$", md)
        meta["title"] = (m.group(1).strip() if m else Path(path).stem)
    return meta, md

# ----------------------------------------------------------------------------- DOCX
def ingest_docx(path, assets_dir=None, keep_assets=True):
    import docx
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph
    d = docx.Document(path)
    meta = {"title": (d.core_properties.title or ""), "author": (d.core_properties.author or ""),
            "format": "docx", "n": len(d.paragraphs)}
    # images (Req 1): extract every image part once; map rId -> markdown ref
    imgmap = {}
    if keep_assets and assets_dir:
        Path(assets_dir).mkdir(parents=True, exist_ok=True)
        n = 0
        for rid, part in list(d.part.related_parts.items()):
            ct = getattr(part, "content_type", "")
            if ct.startswith("image/"):
                n += 1
                ext = ct.split("/")[-1].split("+")[0].replace("jpeg", "jpg").replace("x-emf", "emf")
                fn = f"img-{n}.{ext}"
                try:
                    (Path(assets_dir) / fn).write_bytes(part.blob)
                    imgmap[rid] = f"\n![Figure {n}](assets/{fn})\n"
                except Exception:
                    pass

    def runs_md(p):
        s = ""
        for r in p.runs:
            if keep_assets:
                for bl in r._element.findall(".//" + qn("a:blip")):
                    rid = bl.get(qn("r:embed")) or bl.get(qn("r:link"))
                    if rid in imgmap:
                        s += imgmap[rid]
            t = r.text
            if not t:
                continue
            if r.bold and r.italic: t = f"***{t}***"
            elif r.bold: t = f"**{t}**"
            elif r.italic: t = f"*{t}*"
            s += t
        # Office Math (OMML) -> best-effort inline math (Req 2; preserves the symbols)
        for om in p._p.findall(".//" + qn("m:oMath")):
            mt = "".join((t.text or "") for t in om.findall(".//" + qn("m:t")))
            if mt.strip():
                s += f" ${mt.strip()}$ "
        return s.strip() or p.text.strip()

    def para_md(p):
        st = (p.style.name or "").lower(); txt = runs_md(p)
        if not txt: return ""
        if st.startswith("title") or st.startswith("heading 1"): return f"# {txt}"
        if st.startswith("heading 2"): return f"## {txt}"
        if st.startswith("heading"): return f"### {txt}"
        if "quote" in st: return f"> {txt}"
        if "list" in st or st.startswith("bullet"): return f"- {txt}"
        return txt

    def table_md(tbl):
        rows = []
        for ri, row in enumerate(tbl.rows):
            cells = [c.text.strip().replace("\n", " ").replace("|", r"\|") for c in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
            if ri == 0:
                rows.append("|" + "|".join(["---"] * len(cells)) + "|")
        return "\n" + "\n".join(rows) + "\n"

    out = []
    for child in d.element.body.iterchildren():
        if child.tag == qn("w:p"):
            out.append(para_md(Paragraph(child, d)))
        elif child.tag == qn("w:tbl") and keep_assets:
            out.append(table_md(Table(child, d)))
    md = re.sub(r"\n{3,}", "\n\n", "\n\n".join(o for o in out if o is not None)).strip() + "\n"
    n_h1 = len(re.findall(r"(?m)^# ", md))
    meta["is_book_hint"] = bool(n_h1 >= 4 or wc(md) >= 7000)
    if not meta["title"]:
        m = re.search(r"(?m)^#+ (.+)$", md); meta["title"] = (m.group(1).strip() if m else Path(path).stem)
    return meta, md

# ----------------------------------------------------------------------------- PPTX
def ingest_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    meta = {"title": "", "author": "", "format": "pptx", "n": len(prs.slides)}
    out = []
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        try:
            if slide.shapes.title and slide.shapes.title.text.strip():
                title = slide.shapes.title.text.strip()
        except Exception:
            pass
        out.append(f"\n## Slide {i}" + (f" — {title}" if title else ""))
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            if title and sh == slide.shapes.title:
                continue
            for para in sh.text_frame.paragraphs:
                t = "".join(r.text for r in para.runs).strip() or para.text.strip()
                if t:
                    out.append(("- " if para.level or len(para.runs) else "") + t if t.startswith(("•", "-")) is False else t)
        if i == 1 and title:
            meta["title"] = title
    md = re.sub(r"\n{3,}", "\n\n", "\n".join(out)).strip() + "\n"
    meta["is_book_hint"] = False
    meta["is_slides"] = True
    if not meta["title"]:
        meta["title"] = Path(path).stem
    return meta, md

# ----------------------------------------------------------------------------- HTML / URL
def ingest_html(src, is_url):
    from bs4 import BeautifulSoup
    if is_url:
        import urllib.request
        req = urllib.request.Request(src, headers={"User-Agent": "Mozilla/5.0 (translate-bot)"})
        with urllib.request.urlopen(req, timeout=45) as r:
            raw = r.read().decode("utf-8", "replace")
    else:
        raw = Path(src).read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(raw, "html.parser")
    title = (soup.title.string.strip() if soup.title and soup.title.string else "")
    for tag in soup(["script", "style", "nav", "header", "footer", "aside", "form", "noscript"]):
        tag.decompose()
    main = soup.find("article") or soup.find("main") or soup.body or soup
    out = []
    for el in main.find_all(["h1", "h2", "h3", "h4", "p", "li", "blockquote", "pre"]):
        t = el.get_text(" ", strip=True)
        if not t:
            continue
        nm = el.name
        if nm == "h1": out.append(f"\n# {t}\n")
        elif nm == "h2": out.append(f"\n## {t}\n")
        elif nm in ("h3", "h4"): out.append(f"\n### {t}\n")
        elif nm == "li": out.append(f"- {t}")
        elif nm == "blockquote": out.append(f"> {t}")
        else: out.append(t + "\n")
    md = re.sub(r"\n{3,}", "\n\n", "\n".join(out)).strip() + "\n"
    if not title:
        m = re.search(r"(?m)^# (.+)$", md); title = m.group(1).strip() if m else "Webpage"
    meta = {"title": title, "author": "", "format": "html", "n": wc(md),
            "is_book_hint": len(re.findall(r"(?m)^# ", md)) >= 4}
    return meta, md

# ----------------------------------------------------------------------------- MD / TXT
def ingest_md(path):
    md = Path(path).read_text(encoding="utf-8", errors="replace")
    m = re.search(r"(?m)^#\s+(.+)$", md)
    meta = {"title": (m.group(1).strip() if m else Path(path).stem), "author": "",
            "format": "md", "n": wc(md),
            "is_book_hint": len(re.findall(r"(?m)^# ", md)) >= 4 or wc(md) >= 8000}
    return meta, md.strip() + "\n"

def ingest_txt(path):
    raw = Path(path).read_text(encoding="utf-8", errors="replace")
    paras = [p.strip().replace("\n", " ") for p in re.split(r"\n\s*\n", raw) if p.strip()]
    md = "\n\n".join(paras)
    meta = {"title": Path(path).stem, "author": "", "format": "txt", "n": wc(md),
            "is_book_hint": wc(md) >= 8000}
    return meta, md + "\n"

def ingest_latex(path):
    """Best-effort LaTeX (.tex) -> markdown: translate prose, PRESERVE math + structure."""
    raw = Path(path).read_text(encoding="utf-8", errors="replace")
    m = re.search(r"\\begin\{document\}(.*?)\\end\{document\}", raw, re.S)
    body = m.group(1) if m else raw
    tm = re.search(r"\\title\{(.+?)\}", raw, re.S)
    title = re.sub(r"\\\w+|[{}]", "", tm.group(1)).strip() if tm else Path(path).stem
    body = re.sub(r"(?<!\\)%.*", "", body)                       # strip comments
    body = re.sub(r"\\(label|cite[a-z]*|ref|index)\{[^}]*\}", "", body)
    body = re.sub(r"\\chapter\*?\{(.+?)\}", r"\n# \1\n", body)
    body = re.sub(r"\\section\*?\{(.+?)\}", r"\n## \1\n", body)
    body = re.sub(r"\\subsection\*?\{(.+?)\}", r"\n### \1\n", body)
    body = re.sub(r"\\textbf\{(.+?)\}", r"**\1**", body)
    body = re.sub(r"\\(?:textit|emph)\{(.+?)\}", r"*\1*", body)
    body = re.sub(r"\\includegraphics(?:\[[^\]]*\])?\{(.+?)\}", r"\n![](\1)\n", body)
    body = re.sub(r"\\item\s*", r"\n- ", body)
    body = re.sub(r"\\begin\{(itemize|enumerate|center|figure|table)\}|\\end\{(itemize|enumerate|center|figure|table)\}", "", body)
    body = re.sub(r"\\(maketitle|tableofcontents|centering|hline|toprule|midrule|bottomrule)\b", "", body)
    body = re.sub(r"[ \t]+\n", "\n", body)
    body = re.sub(r"\n{3,}", "\n\n", body).strip()
    md = f"# {title}\n\n{body}\n"
    meta = {"title": title, "author": "", "format": "latex", "n": wc(md),
            "is_book_hint": len(re.findall(r"(?m)^# ", md)) >= 3 or bool(re.search(r"\\chapter", raw))}
    return meta, md

def ingest(cfg, assets_dir=None, keep_assets=True):
    fmt, src = cfg["srcFormat"], cfg["src"]
    if fmt == "pdf": return ingest_pdf(src, assets_dir, keep_assets, cfg.get("math", "auto"))
    if fmt == "docx": return ingest_docx(src, assets_dir, keep_assets)
    if fmt == "pptx": return ingest_pptx(src)
    if fmt == "html": return ingest_html(src, cfg["isUrl"])
    if fmt == "md": return ingest_md(src)
    if fmt == "txt": return ingest_txt(src)
    if fmt == "latex": return ingest_latex(src)
    raise SystemExit(f"unsupported source format: {fmt}")

# ----------------------------------------------------------------------------- split
def split_units(md, meta, mode):
    """Return list of units [{id,title,is_first,part,md,words}]."""
    lines = md.splitlines()
    # find H1 boundaries (chapters). If none, the whole doc is one chapter.
    h1 = [i for i, l in enumerate(lines) if l.startswith("# ")]
    chapters = []
    if mode in ("book",) and len(h1) >= 2:
        bounds = h1 + [len(lines)]
        for k in range(len(h1)):
            seg = lines[bounds[k]:bounds[k + 1]]
            title = seg[0][2:].strip()
            chapters.append((title, "\n".join(seg).strip()))
    else:
        # doc / slides / book-without-H1: one chapter = whole doc (title from meta)
        title = meta["title"]
        body = md if md.lstrip().startswith("# ") else f"# {title}\n\n{md}"
        chapters.append((title, body.strip()))

    # Merge tiny fragments (title-page echoes, blank dividers, single-line front
    # matter) into a neighbor so generic book detection doesn't emit junk units.
    if len(chapters) > 1:
        merged = []
        for title, body in chapters:
            tiny = wc(re.sub(r"(?m)^#.*$", "", body)) < 45
            if merged and tiny:
                pt, pb = merged[-1]
                merged[-1] = (pt, pb + "\n\n" + re.sub(r"^#\s.*\n?", "", body.lstrip(), count=1).strip())
            elif merged and wc(re.sub(r"(?m)^#.*$", "", merged[-1][1])) < 45:
                # previous opener was itself tiny -> adopt THIS real title, keep both bodies
                pt, pb = merged[-1]
                merged[-1] = (title, pb + "\n\n" + body)
            else:
                merged.append((title, body))
        chapters = merged

    units = []
    for ci, (title, cmd) in enumerate(chapters, 1):
        words = wc(cmd)
        nparts = max(1, math.ceil(words / CAP))
        clines = cmd.splitlines()
        if nparts == 1:
            units.append({"num": f"{ci:02d}", "title": title, "is_first": True, "part": "",
                          "md": cmd, "words": words})
            continue
        sec = [i for i, l in enumerate(clines) if l.startswith("## ")]
        cum, run = [], 0
        for l in clines:
            run += wc(l); cum.append(run)
        cuts = []
        for kk in range(1, nparts):
            tgt = words * kk / nparts
            if sec:
                best = min(sec, key=lambda i: abs(cum[i] - tgt))
                if 0 < best < len(clines) and best not in cuts:
                    cuts.append(best)
        cuts = sorted(set(cuts)); bnd = [0] + cuts + [len(clines)]
        letters = "abcdefghijklmnop"
        for pi in range(len(bnd) - 1):
            chunk = "\n".join(clines[bnd[pi]:bnd[pi + 1]]).strip()
            if pi > 0:
                chunk = re.sub(r"^#\s+.*\n?", "", chunk)  # continuation: drop stray H1
                chunk = f"{chunk}"
            units.append({"num": f"{ci:02d}{letters[pi]}", "title": title,
                          "is_first": pi == 0, "part": letters[pi], "md": chunk, "words": wc(chunk)})
    return units

# ----------------------------------------------------------------------------- brief
def write_brief(workdir, cfg):
    to = cfg["to"]; lang = to["name"]; endo = to.get("endonym", lang)
    frm = (cfg.get("from") or {}).get("name", "the source language")
    gloss = ""
    if cfg.get("glossary") and Path(cfg["glossary"]).exists():
        gloss = "\n\n## Project glossary (use consistently)\n\n" + Path(cfg["glossary"]).read_text(encoding="utf-8")
    # v0.3: math rule branches on --math. Default (auto/preserve) = preserve clean LaTeX AND
    # reconstruct garbled math (PDF text-layer extraction mangles equations); crop = display
    # equations arrive as image crops (don't reconstruct); off = leave math untouched.
    math_mode = cfg.get("math", "auto")
    if math_mode == "off":
        rule4 = (
            "4. **MATH & FORMULAS — leave EXACTLY as-is.** Do not translate, alter, or "
            "reconstruct anything that looks like a formula or symbol. Translate only the "
            "surrounding prose.")
    elif math_mode == "crop":
        rule4 = (
            "4. **MATH & FORMULAS.** Display equations are supplied as IMAGE CROPS — "
            "`![Eq. N](assets/eq-...png)` lines: keep the `(assets/...)` path EXACTLY, do "
            "not reconstruct them. Any already-clean inline `$...$` → copy byte-for-byte. "
            "Translate only the surrounding prose.")
    else:
        rule4 = (
            "4. **MATH & FORMULAS — native LaTeX (preserve clean, reconstruct garbled):**\n"
            "   - Anything already in clean LaTeX — `$...$`, `$$...$$`, `\\(...\\)`, `\\[...\\]`, or a\n"
            "     math environment (`aligned`, `cases`, `matrix`, …) — **copy BYTE-FOR-BYTE** (every\n"
            "     symbol, sub/superscript, `\\command`, brace). Keep `$$` blocks on their own lines.\n"
            "   - This source may be a PDF whose math extracted as **scrambled text** (broken glyphs\n"
            "     like `q⇤(a) .= E[Rt|At=a]`, multi-line equations shattered into fragments, subscripts\n"
            "     mis-tagged as `<sup>`). When you see mangled math, **RECONSTRUCT the intended formula\n"
            "     into proper LaTeX** using the surrounding prose for context: inline → `$...$`, display\n"
            "     → `$$...$$` on its own lines (multi-line derivations → `\\begin{aligned}…\\end{aligned}`,\n"
            "     numbered → trailing `\\qquad (N.M)`). Use `\\doteq`, `\\varepsilon`, `\\mathbf{1}` for the\n"
            "     indicator, `\\operatorname*{arg\\,max}`, etc. Reconstruct the MATH only — never invent content.\n"
            "   - **Algorithm/pseudocode boxes** → a bold heading line, then a TOP-LEVEL\n"
            "     `$$\\begin{aligned}…\\end{aligned}$$` block (use `\\begin{cases}` for branches). Do NOT\n"
            "     wrap pseudocode in a `>` blockquote.\n"
            "   - **`\\text{…}` spans inside math** (pseudocode words, worded labels like NewEstimate):\n"
            "     translate the natural-language words into " + lang + ", keep all symbols/commands/braces.\n"
            "   - Do NOT translate variable names, operators, or function names inside math.")
    brief = f"""# Translator brief — into {lang} ({endo})

You are a **world-class literary/non-fiction translator into {lang}** — the caliber a
top publishing house hires. Translate the given unit from {frm} into **{lang}**. The
result must be **faithful AND beautiful** — it should read as if originally written in
{lang}, never as a machine translation.

## Iron rules
1. **Faithful, not literal.** Convey meaning, nuance, tone, and rhythm. Recast syntax
   into natural {lang} word order; split/merge sentences where {lang} flows better.
   Add nothing; omit nothing of substance.
2. **Native fluency.** Idioms → {lang} equivalents, never calques. Read it in your head:
   it must sound like a gifted {lang} writer. Maintain a consistent register and a
   consistent voice for each speaker.
3. **Preserve every bit of structure (markdown):**
   - `#`, `##`, `###` headings → translate the text, keep the level.
   - `>` blockquotes → translate, keep as blockquote.
   - `**bold**`, `*italic*`, lists (`-`, `1.`), and tables → keep the markup (translate cell text).
   - **`<sup>N</sup>` footnote markers → keep VERBATIM**, attached to the preceding word.
   - **Image/figure references `![alt](path)` → keep the `(path)` EXACTLY; translate only the
     alt/caption text inside `[...]`.** Never alter, drop, or reorder a figure reference.
{rule4}
5. **Proper nouns & terms:** keep personal names, brands, and product names in their
   original form (e.g. Tesla, SpaceX). Localize only well-established place names per
   {lang} convention. Keep numbers, units, dates, and currency as-is.
6. **Consistency:** translate recurring key terms the same way throughout.{gloss}

## This unit
{("This is the FIRST part of the chapter — begin your output with the H1 line you are given (translated), then the body."
  if True else "")}
- If you are told this is a CONTINUATION part, do NOT repeat the chapter title; begin
  directly with the first translated section/paragraph.

## Output
Return **ONLY** the translated {lang} markdown — no preamble, no code fences, no source
text, no notes. Preserve every `<sup>N</sup>`. Make it sing.
"""
    Path(workdir, "brief.md").write_text(brief, encoding="utf-8")

# ----------------------------------------------------------------------------- plan
def cmd_plan(workdir):
    wd = Path(workdir)
    cfg = json.loads((wd / "config.json").read_text())
    assets_dir = wd / "assets"
    meta, md = ingest(cfg, str(assets_dir), cfg.get("keepAssets", True))
    (wd / "normalized.md").write_text(md, encoding="utf-8")
    meta["assets"] = len(list(assets_dir.glob("*"))) if assets_dir.exists() else 0
    mode = cfg["mode"]
    if mode == "auto":
        mode = "slides" if meta.get("is_slides") else ("book" if meta.get("is_book_hint") else "doc")
    units = split_units(md, meta, mode)
    usrc = wd / "units"; usrc.mkdir(exist_ok=True)
    for u in units:
        head = f"# {u['title']}\n\n" if u["is_first"] and not u["md"].lstrip().startswith("# ") else ""
        body = u["md"] if not head else head + re.sub(r"^#\s+.*\n?", "", u["md"].lstrip(), count=1)
        (usrc / f"{u['num']}-{slugify(u['title'])}.src.md").write_text(
            (body if u["is_first"] else u["md"]).strip() + "\n", encoding="utf-8")
        u["slug"] = slugify(u["title"])
        u["src_file"] = str((usrc / f"{u['num']}-{u['slug']}.src.md").resolve())
        u["out_file"] = str((wd / f"units_{cfg['to']['code']}" / f"{u['num']}-{u['slug']}.tr.md").resolve())
    (wd / f"units_{cfg['to']['code']}").mkdir(exist_ok=True)
    write_brief(wd, cfg)
    total = sum(u["words"] for u in units)
    plan = {"mode": mode, "meta": meta, "units": units, "total_words": total,
            "cost_estimate_usd": round(total / 1000 * 0.06, 2), "outputs": cfg["out"],
            "to": cfg["to"], "style": cfg["style"], "brief": str((wd / "brief.md").resolve())}
    (wd / "plan.json").write_text(json.dumps(plan, ensure_ascii=False, indent=2), encoding="utf-8")
    (wd / "meta.json").write_text(json.dumps(meta, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps({"mode": mode, "title": meta["title"], "units": len(units),
                      "total_words": total, "cost_estimate_usd": plan["cost_estimate_usd"],
                      "outputs": cfg["out"], "to": cfg["to"]["name"], "workdir": str(wd.resolve())},
                     ensure_ascii=False))

# ----------------------------------------------------------------------------- assemble
def cmd_assemble(workdir):
    wd = Path(workdir)
    cfg = json.loads((wd / "config.json").read_text())
    plan = json.loads((wd / "plan.json").read_text())
    meta = plan["meta"]; code = cfg["to"]["code"]
    udir = wd / f"units_{code}"
    # group split parts (NNa,NNb) back into chapters by the 2-digit chapter number
    bychap = collections.OrderedDict()
    for u in plan["units"]:
        chap = u["num"][:2]
        f = udir / f"{u['num']}-{u['slug']}.tr.md"
        txt = f.read_text(encoding="utf-8").strip() if f.exists() else ""
        bychap.setdefault(chap, {"title": u["title"], "parts": []})["parts"].append(txt)
    blocks = []
    for chap, c in bychap.items():
        md = "\n\n".join(p for p in c["parts"] if p)
        if not md.lstrip().startswith("# "):
            md = f"# {c['title']}\n\n{md}"
        blocks.append({"type": "chapter", "num": chap, "id": f"ch{chap}",
                       "title": c["title"], "md": md})
    out = {"meta": {"title_src": meta.get("title", ""), "author": meta.get("author", ""),
                    "format": meta.get("format"), "to": cfg["to"], "style": cfg["style"],
                    "title_translated": cfg.get("titleTranslated") or meta.get("title", ""),
                    "subtitle": cfg.get("subtitle", "")},
           "mode": plan["mode"], "blocks": blocks}
    (wd / "book_blocks.json").write_text(json.dumps(out, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps({"chapters": len(blocks), "mode": plan["mode"]}))

if __name__ == "__main__":
    cmd = sys.argv[1] if len(sys.argv) > 1 else ""
    wd = sys.argv[2] if len(sys.argv) > 2 else "."
    if cmd == "plan": cmd_plan(wd)
    elif cmd == "assemble": cmd_assemble(wd)
    else: raise SystemExit("usage: engine.py {plan|assemble} <workdir>")
