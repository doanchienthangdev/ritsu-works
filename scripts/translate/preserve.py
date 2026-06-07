#!/usr/bin/env python3
"""preserve.py — format-preserving (in-place) translation for /translate --preserve (Req 3).
Translates the TEXT of a docx/pptx/pdf in place, leaving structure, styling, images, tables,
and layout identical to the original. Two subcommands mirroring the reflow pipeline so the
SAME Workflow translates the units:
  preserve.py plan  <workdir>   extract ordered text segments -> units/NN.src.md (⟦S{id}⟧ lines)
  preserve.py build <workdir>   reinsert translated segments into a copy of the original -> output
"""
import sys, re, json, shutil
from pathlib import Path

SEG = re.compile(r"⟦S(\d+)⟧[ \t]*(.*)")
CAP_SEG = 140          # segments per translation unit
SKIP = re.compile(r"^[\s\d.,:;%+\-–—/()\[\]<>=*#©®™|]*$")  # numeric/punct-only → don't translate

def translatable(t):
    return bool(t) and len(t.strip()) >= 2 and not SKIP.match(t.strip())

# ---------------------------------------------------------------- segment iterators (order-stable)
def docx_paras(doc):
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph
    def walk(elm):
        for child in elm.iterchildren():
            if child.tag == qn("w:p"):
                yield Paragraph(child, doc)
            elif child.tag == qn("w:tbl"):
                for row in Table(child, doc).rows:
                    for cell in row.cells:
                        for p in cell.paragraphs:
                            yield p
    yield from walk(doc.element.body)

def pptx_paras(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    yield para

def set_runs_text(runs_holder, text):
    """Set the holder's first text run to `text`, blank the rest (keeps formatting + drawings)."""
    runs = [r for r in runs_holder.runs if (r.text or "").strip()]
    if not runs:
        if hasattr(runs_holder, "add_run"):
            runs_holder.add_run(text)
        elif hasattr(runs_holder, "text"):
            runs_holder.text = text
        return
    runs[0].text = text
    for r in runs[1:]:
        r.text = ""

# ---------------------------------------------------------------- extract / reinsert per format
def iter_segments(src, fmt):
    """Yield (holder, text) for each translatable unit, in stable order."""
    if fmt == "docx":
        import docx
        for p in docx_paras(docx.Document(src)):
            yield p, (p.text or "")
    elif fmt == "pptx":
        from pptx import Presentation
        for para in pptx_paras(Presentation(src)):
            yield para, ("".join(r.text for r in para.runs) or "")
    else:
        raise SystemExit(f"--preserve does not support source format {fmt}")

def collect(src, fmt):
    segs = []
    for _, text in iter_segments(src, fmt):
        if translatable(text):
            segs.append(text.strip())
        else:
            segs.append(None)  # keep slot for stable ids; not translated
    return segs

# ---------------------------------------------------------------- plan
def cmd_plan(wd):
    wd = Path(wd); cfg = json.loads((wd / "config.json").read_text())
    src, fmt = cfg["src"], cfg["srcFormat"]
    segs = collect(src, fmt)
    (wd / "segments.json").write_text(json.dumps(segs, ensure_ascii=False), encoding="utf-8")
    ids = [i for i, s in enumerate(segs) if s is not None]
    usrc = wd / "units"; usrc.mkdir(exist_ok=True)
    (wd / f"units_{cfg['to']['code']}").mkdir(exist_ok=True)
    units = []
    for k in range(0, len(ids), CAP_SEG):
        batch = ids[k:k + CAP_SEG]
        num = f"{k // CAP_SEG + 1:02d}"
        lines = [f"⟦S{i}⟧ {segs[i]}" for i in batch]
        (usrc / f"{num}-seg.src.md").write_text("\n\n".join(lines) + "\n", encoding="utf-8")
        units.append({"num": num, "title": "segments", "is_first": True, "part": "", "slug": "seg",
                      "words": sum(len(segs[i].split()) for i in batch),
                      "src_file": str((usrc / f"{num}-seg.src.md").resolve()),
                      "out_file": str((wd / f"units_{cfg['to']['code']}" / f"{num}-seg.tr.md").resolve())})
    write_preserve_brief(wd, cfg)
    total = sum(u["words"] for u in units)
    plan = {"mode": "preserve", "meta": {"title": Path(src).stem, "format": fmt, "segments": len(ids)},
            "units": units, "total_words": total, "cost_estimate_usd": round(total / 1000 * 0.06, 2),
            "outputs": [fmt], "to": cfg["to"], "style": cfg["style"], "brief": str((wd / "brief.md").resolve())}
    (wd / "plan.json").write_text(json.dumps(plan, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps({"mode": "preserve", "title": Path(src).stem, "units": len(units),
                      "total_words": total, "cost_estimate_usd": plan["cost_estimate_usd"],
                      "outputs": [fmt], "to": cfg["to"]["name"], "workdir": str(wd.resolve()),
                      "segments": len(ids)}, ensure_ascii=False))

def write_preserve_brief(wd, cfg):
    to = cfg["to"]; lang = to["name"]
    (wd / "brief.md").write_text(f"""# Translator brief — into {lang} (format-preserving / in-place)

You are a world-class translator into {lang}. The input is a list of text segments, each on
its own line beginning with a marker like `⟦S12⟧`. Translate ONLY the text after each marker
into natural, faithful {lang}.

IRON RULES (a broken segment map corrupts the document):
- Keep EVERY `⟦S{{n}}⟧` marker EXACTLY (same number), one segment per line, SAME ORDER.
- Do NOT merge, split, add, drop, or reorder segments. One input segment → one output line.
- If a segment is a proper noun, code, a number, or already in {lang}, keep it as-is.
- **Preserve math VERBATIM**: anything in `$...$`, `\\(...\\)`, `\\[...\\]` — copy byte-for-byte.
- Keep URLs, file paths, and product/brand names unchanged.
- A segment is usually one paragraph or table cell; translate it as natural {lang} prose.

OUTPUT: return ONLY the translated segments, each line `⟦S{{n}}⟧ <translation>`, blank line
between segments, no preamble, no notes.
""", encoding="utf-8")

# ---------------------------------------------------------------- build (reinsert)
def load_translations(wd, code):
    tr = {}
    udir = Path(wd) / f"units_{code}"
    for f in sorted(udir.glob("*.tr.md")):
        for m in SEG.finditer(f.read_text(encoding="utf-8")):
            tr[int(m.group(1))] = m.group(2).strip()
    return tr

def cmd_build(wd):
    wd = Path(wd); cfg = json.loads((wd / "config.json").read_text())
    src, fmt = cfg["src"], cfg["srcFormat"]
    code = cfg["to"]["code"]
    tr = load_translations(wd, code)
    out_path = str(Path(cfg["outDir"]) / f"{cfg['name']}.{code}.{fmt}")
    Path(cfg["outDir"]).mkdir(parents=True, exist_ok=True)

    # ids are FULL paragraph indices (matching plan's enumerate), so enumerate ALL here too
    if fmt == "docx":
        import docx
        doc = docx.Document(src)
        for idx, p in enumerate(docx_paras(doc)):
            if translatable(p.text or "") and tr.get(idx):
                set_runs_text(p, tr[idx])
        doc.save(out_path)
    elif fmt == "pptx":
        from pptx import Presentation
        prs = Presentation(src)
        for idx, para in enumerate(pptx_paras(prs)):
            t = "".join(r.text for r in para.runs)
            if translatable(t) and tr.get(idx):
                set_runs_text(para, tr[idx])
        prs.save(out_path)
    else:
        raise SystemExit(f"--preserve build does not support {fmt}")
    print(json.dumps({"fmt": fmt, "out": out_path, "segments_translated": len(tr)}))

if __name__ == "__main__":
    cmd = sys.argv[1] if len(sys.argv) > 1 else ""
    wd = sys.argv[2] if len(sys.argv) > 2 else "."
    if cmd == "plan": cmd_plan(wd)
    elif cmd == "build": cmd_build(wd)
    else: raise SystemExit("usage: preserve.py {plan|build} <workdir>")
