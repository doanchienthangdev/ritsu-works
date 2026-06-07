#!/usr/bin/env python3
"""render.py — render ONE output format from a translated book_blocks.json.
  python render.py <workdir> <fmt> [--fonts=DIR] [--repo-root=DIR]
fmt: pdf (WeasyPrint host) | epub | docx | pptx | md (anaconda host).
Design palette comes from the resolved style; typography = bundled Source Serif 4 + Inter.
"""
import sys, os, re, json, html as H
from pathlib import Path

HERE = Path(__file__).parent
sys.path.insert(0, str(HERE))
import design  # noqa
import latexout  # noqa

def opts(argv):
    o = {}
    for a in argv:
        if a.startswith("--"):
            k, _, v = a[2:].partition("=")
            o[k] = v or True
    return o

def md_html(text):
    import markdown
    text = re.sub(r"^#\s+.*\n", "", text.lstrip(), count=1)  # strip chapter H1 (rendered separately)
    return markdown.markdown(text, extensions=["extra", "sane_lists", "smarty"]).replace("<br>", "<br/>")

def load(wd):
    data = json.loads((Path(wd) / "book_blocks.json").read_text(encoding="utf-8"))
    return data["meta"], data["blocks"], data.get("mode", "book")

def endonym(meta): return meta["to"].get("endonym", meta["to"]["name"])

# ----------------------------------------------------------------------------- PDF
def render_pdf(wd, out_path, fonts_dir, repo_root):
    from weasyprint import HTML, CSS
    from weasyprint.text.fonts import FontConfiguration
    meta, blocks, mode = load(wd)
    tokens, _ = design.load_tokens(meta["style"], repo_root)
    title = meta.get("title_translated") or meta.get("title_src") or "—"
    author = meta.get("author", "")
    parts = []
    toc = []
    for b in blocks:
        body = md_html(b["md"])
        parts.append(f'<section class="chapter" id="{b["id"]}"><div class="ch-ornament">✦</div>'
                     f'<h1 class="chapter-title">{H.escape(b["title"])}</h1>'
                     f'<div class="ch-body">{body}</div></section>')
        toc.append(f'<li class="toc-ch"><a href="#{b["id"]}">{H.escape(b["title"])}</a></li>')
    cover = (f'<section class="cover"><div class="cover-frame">'
             + (f'<div class="cover-kicker">{H.escape(author)}</div>' if author else '')
             + f'<h1 class="cover-title">{H.escape(title)}</h1>'
             + (f'<div class="cover-sub">{H.escape(meta.get("subtitle",""))}</div>' if meta.get("subtitle") else '')
             + f'<div class="cover-mark">✦</div>'
             + (f'<div class="cover-author">{H.escape(author)}</div>' if author else '')
             + f'<div class="cover-foot">Bản dịch · {H.escape(meta["to"]["name"])} ({H.escape(endonym(meta))})</div>'
             + '</div></section>')
    colophon = (f'<section class="colophon"><div class="colo-orig">Nguyên tác</div>'
                f'<div class="colo-origtitle">{H.escape(meta.get("title_src") or title)}</div>'
                + (f'<div class="colo-credits">{H.escape(author)}</div>' if author else '')
                + f'<div class="colo-note">Bản dịch sang {H.escape(meta["to"]["name"])} được tạo cho mục đích đọc/'
                f'nghiên cứu cá nhân. Mọi quyền đối với nguyên tác thuộc về tác giả và nhà xuất bản gốc.</div>'
                f'<div class="colo-design">Thiết kế theo hệ thống thị giác <b>{H.escape(tokens.get("_name","Claude"))}</b> '
                f'· Source Serif 4 &amp; Inter · /translate</div></section>')
    tocsec = '<section class="toc"><h1 class="toc-h">Mục Lục</h1><ul>' + "".join(toc) + "</ul></section>"
    show_front = (mode != "doc") or len(blocks) > 1
    html = "<!doctype html><html lang='%s'><head><meta charset='utf-8'></head><body>%s</body></html>" % (
        meta["to"]["code"], (cover + colophon + tocsec if show_front else "") + "".join(parts))
    fc = FontConfiguration()
    HTML(string=html, base_url=str(Path(wd).resolve())).write_pdf(
        out_path, stylesheets=[CSS(string=design.book_css(tokens, fonts_dir), font_config=fc)], font_config=fc)

# ----------------------------------------------------------------------------- cover image (PIL) for EPUB
def cover_png(meta, tokens, fonts_dir, dest):
    from PIL import Image, ImageDraw, ImageFont
    W, Hh = 1200, 1856
    img = Image.new("RGB", (W, Hh), tokens["background"]); d = ImageDraw.Draw(img)
    d.rectangle([70, 90, W - 70, Hh - 90], outline=tokens["primary"], width=4)
    serif = lambda sz: ImageFont.truetype(str(Path(fonts_dir) / "SourceSerif4-Bold.ttf"), sz)
    serI = lambda sz: ImageFont.truetype(str(Path(fonts_dir) / "SourceSerif4-Italic.ttf"), sz)
    sans = lambda sz: ImageFont.truetype(str(Path(fonts_dir) / "Inter-SemiBold.ttf"), sz)
    title = meta.get("title_translated") or meta.get("title_src") or "—"
    def wrap(text, font, maxw):
        words, lines, cur = text.split(), [], ""
        for w in words:
            t = (cur + " " + w).strip()
            if d.textlength(t, font=font) <= maxw: cur = t
            else: lines.append(cur); cur = w
        if cur: lines.append(cur)
        return lines
    y = 360
    for ln in wrap(title, serif(96), W - 320):
        wln = d.textlength(ln, font=serif(96)); d.text(((W - wln) / 2, y), ln, font=serif(96), fill=tokens["foreground"]); y += 118
    if meta.get("author"):
        a = meta["author"]; wa = d.textlength(a, font=sans(40)); d.text(((W - wa) / 2, Hh - 360), a, font=sans(40), fill=tokens["foreground"])
    foot = f'Bản dịch · {meta["to"]["name"]}'
    wf = d.textlength(foot, font=serI(34)); d.text(((W - wf) / 2, Hh - 290), foot, font=serI(34), fill=tokens["mutedForeground"])
    d.text((W / 2 - 12, Hh / 2 + 120), "✦", font=serif(48), fill=tokens["primary"])
    img.save(dest, "PNG")

# ----------------------------------------------------------------------------- EPUB
def render_epub(wd, out_path, fonts_dir, repo_root):
    from ebooklib import epub
    meta, blocks, mode = load(wd)
    tokens, _ = design.load_tokens(meta["style"], repo_root)
    title = meta.get("title_translated") or meta.get("title_src") or "—"
    book = epub.EpubBook(); book.set_identifier("translate-" + re.sub(r"\W+", "-", title.lower())[:40])
    book.set_title(title); book.set_language(meta["to"]["code"])
    if meta.get("author"): book.add_author(meta["author"])
    covimg = Path(wd) / "cover.png"; cover_png(meta, tokens, fonts_dir, covimg)
    book.set_cover("cover.png", covimg.read_bytes())
    css = design.epub_css(tokens, "../fonts")
    for fn in ["SourceSerif4-Regular.ttf", "SourceSerif4-SemiBold.ttf", "SourceSerif4-Bold.ttf",
               "SourceSerif4-Italic.ttf", "Inter-Regular.ttf", "Inter-SemiBold.ttf", "Inter-Bold.ttf"]:
        book.add_item(epub.EpubItem(uid="f_" + fn, file_name="fonts/" + fn,
                                    content=(Path(fonts_dir) / fn).read_bytes(), media_type="font/ttf"))
    style = epub.EpubItem(uid="style", file_name="style/book.css", media_type="text/css", content=css)
    book.add_item(style)
    def page(uid, fn, ttl, body, cls=""):
        it = epub.EpubHtml(uid=uid, file_name=fn, title=ttl, lang=meta["to"]["code"]); it.add_item(style)
        it.content = (f"<html xmlns='http://www.w3.org/1999/xhtml'><head><title>{H.escape(ttl)}</title>"
                      f"<link rel='stylesheet' href='../style/book.css'/></head><body class='{cls}'>{body or '<p></p>'}</body></html>")
        book.add_item(it); return it
    tp = page("tp", "text/000-title.xhtml", title,
              f"<div class='tp'><h1>{H.escape(title)}</h1>"
              + (f"<p class='s'>{H.escape(meta.get('subtitle',''))}</p>" if meta.get('subtitle') else "")
              + "<div class='orn'>✦</div>"
              + (f"<div class='a'>{H.escape(meta.get('author',''))}</div>" if meta.get('author') else "")
              + f"<p class='s' style='font-size:.8em'>Bản dịch · {H.escape(meta['to']['name'])}</p></div>")
    spine = ["nav", tp]; toc = []
    for i, b in enumerate(blocks, 1):
        body = f"<div class='orn'>✦</div><h1 class='ct'>{H.escape(b['title'])}</h1><div>{md_html(b['md'])}</div>"
        it = page(b["id"], f"text/{i:03d}-{b['id']}.xhtml", b["title"], body)
        spine.append(it); toc.append(epub.Link(f"text/{i:03d}-{b['id']}.xhtml", b["title"], b["id"]))
    book.toc = tuple(toc); book.add_item(epub.EpubNcx()); book.add_item(epub.EpubNav()); book.spine = spine
    epub.write_epub(out_path, book)

# ----------------------------------------------------------------------------- DOCX
def render_docx(wd, out_path, fonts_dir, repo_root):
    import docx
    from docx.shared import Pt, RGBColor, Inches
    meta, blocks, mode = load(wd)
    tokens, _ = design.load_tokens(meta["style"], repo_root)
    deep = RGBColor.from_string(tokens["primaryDeep"].lstrip("#"))
    ink = RGBColor.from_string(tokens["foreground"].lstrip("#"))
    doc = docx.Document()
    n = doc.styles["Normal"]; n.font.name = "Georgia"; n.font.size = Pt(11); n.font.color.rgb = ink
    title = meta.get("title_translated") or meta.get("title_src") or "—"
    t = doc.add_paragraph(); r = t.add_run(title); r.bold = True; r.font.size = Pt(26); r.font.color.rgb = ink
    if meta.get("author"):
        a = doc.add_paragraph(); ar = a.add_run(meta["author"]); ar.font.size = Pt(13); ar.font.color.rgb = deep
    doc.add_paragraph(f"Bản dịch · {meta['to']['name']}").runs[0].italic = True
    doc.add_page_break()
    inline = re.compile(r"(\*\*.+?\*\*|\*.+?\*|<sup>\d+</sup>)")
    def emit(p_el, text):
        for tok in inline.split(text):
            if not tok: continue
            if tok.startswith("**") and tok.endswith("**"): rr = p_el.add_run(tok[2:-2]); rr.bold = True
            elif tok.startswith("*") and tok.endswith("*"): rr = p_el.add_run(tok[1:-1]); rr.italic = True
            elif tok.startswith("<sup>"): rr = p_el.add_run(re.sub("<.*?>", "", tok)); rr.font.superscript = True; rr.font.color.rgb = deep
            else: p_el.add_run(tok)
    for b in blocks:
        h = doc.add_paragraph(); hr = h.add_run(b["title"]); hr.bold = True; hr.font.size = Pt(20); hr.font.color.rgb = ink
        for ln in b["md"].splitlines():
            ln = ln.rstrip()
            if not ln or ln.startswith("# "): continue
            if ln.startswith("## "):
                p = doc.add_paragraph(); rr = p.add_run(ln[3:].upper()); rr.bold = True; rr.font.size = Pt(11); rr.font.color.rgb = deep
            elif ln.startswith("### "):
                p = doc.add_paragraph(); rr = p.add_run(ln[4:]); rr.bold = True; rr.font.size = Pt(13)
            elif ln.startswith("> "):
                p = doc.add_paragraph(); p.paragraph_format.left_indent = Inches(0.4); emit(p, ln[2:]);
                for rr in p.runs: rr.italic = True
            elif ln.startswith("- "):
                p = doc.add_paragraph(style="List Bullet"); emit(p, ln[2:])
            else:
                p = doc.add_paragraph(); emit(p, ln)
        doc.add_page_break()
    doc.save(out_path)

# ----------------------------------------------------------------------------- PPTX
def render_pptx(wd, out_path, fonts_dir, repo_root):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    meta, blocks, mode = load(wd)
    tokens, _ = design.load_tokens(meta["style"], repo_root)
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    bg = tokens["background"].lstrip("#"); ink = tokens["foreground"].lstrip("#"); deep = tokens["primaryDeep"].lstrip("#")
    blank = prs.slide_layouts[6]
    def paint(slide):
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = RGBColor.from_string(bg)
    def textbox(slide, x, y, w, h):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tb.text_frame.word_wrap = True; return tb.text_frame
    # title slide
    s = prs.slides.add_slide(blank); paint(s)
    tf = textbox(s, 1, 2.6, 11.3, 2.5); p = tf.paragraphs[0]
    r = p.add_run(); r.text = meta.get("title_translated") or meta.get("title_src") or "—"
    r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = RGBColor.from_string(ink)
    if meta.get("author"):
        p2 = tf.add_paragraph(); rr = p2.add_run(); rr.text = meta["author"]; rr.font.size = Pt(20); rr.font.color.rgb = RGBColor.from_string(deep)
    for b in blocks:
        # one content slide per H2 section (or the whole chapter if none)
        secs = re.split(r"(?m)^##\s+", b["md"])
        chunks = [(b["title"], secs[0])] if len(secs) == 1 else \
                 [(b["title"], "")] + [(s.split("\n", 1)[0].strip(), s.split("\n", 1)[1] if "\n" in s else "") for s in secs[1:]]
        for stitle, sbody in chunks:
            sl = prs.slides.add_slide(blank); paint(sl)
            tt = textbox(sl, 0.7, 0.5, 12, 1.2); tp = tt.paragraphs[0]; tr = tp.add_run(); tr.text = stitle[:120]
            tr.font.size = Pt(28); tr.font.bold = True; tr.font.color.rgb = RGBColor.from_string(ink)
            body = textbox(sl, 0.9, 1.9, 11.5, 5.0)
            lines = [re.sub(r"<.*?>|\*\*|\*|^>\s?|^-\s?", "", ln).strip()
                     for ln in sbody.splitlines() if ln.strip() and not ln.startswith("#")]
            first = True
            for ln in lines[:12]:
                p = body.paragraphs[0] if first else body.add_paragraph(); first = False
                rr = p.add_run(); rr.text = ("• " + ln) if len(ln) < 200 else ("• " + ln[:200] + "…")
                rr.font.size = Pt(16); rr.font.color.rgb = RGBColor.from_string(ink)
    prs.save(out_path)

# ----------------------------------------------------------------------------- MD
def render_md(wd, out_path, fonts_dir, repo_root):
    meta, blocks, mode = load(wd)
    title = meta.get("title_translated") or meta.get("title_src") or "—"
    out = [f"# {title}", ""]
    if meta.get("author"): out += [f"**{meta['author']}**", ""]
    out += [f"*Bản dịch · {meta['to']['name']} ({endonym(meta)})*", "", "---", ""]
    for b in blocks:
        md = b["md"] if b["md"].lstrip().startswith("# ") else f"# {b['title']}\n\n{b['md']}"
        out += [re.sub(r"^#\s", "## ", md.lstrip(), count=1), ""]
    Path(out_path).write_text("\n".join(out), encoding="utf-8")

# ----------------------------------------------------------------------------- LaTeX / pdf-latex
def _latex_data(wd, repo_root):
    data = json.loads((Path(wd) / "book_blocks.json").read_text(encoding="utf-8"))
    tokens, _ = design.load_tokens(data["meta"]["style"], repo_root)
    return data, tokens

def _stage(src_fonts, src_assets, dest):
    """Copy the bundled TTFs (+ assets if any) into dest; return (fonts_subdir, assets_subdir|None)."""
    import shutil
    dest = Path(dest)
    fdir = dest / "fonts"; fdir.mkdir(parents=True, exist_ok=True)
    for ttf in Path(src_fonts).glob("*.ttf"):
        d = fdir / ttf.name
        if not d.exists():
            shutil.copy2(ttf, d)
    adir = None
    if src_assets and Path(src_assets).exists() and any(Path(src_assets).iterdir()):
        adir = dest / "assets"
        if not adir.exists():
            shutil.copytree(src_assets, adir)
    return fdir, adir

def render_latex(wd, out_path, fonts_dir, repo_root):
    # portable .tex: stage fonts + assets beside it, reference relatively
    data, tokens = _latex_data(wd, repo_root)
    out = Path(out_path); stem = out.stem  # e.g. gd.vi
    support = out.parent / f"{stem}.support"
    fdir, adir = _stage(fonts_dir, Path(wd) / "assets", support)
    rel_f = f"{stem}.support/fonts"
    rel_a = f"{stem}.support/assets" if adir else ""
    tex = latexout.to_latex(data, rel_f, rel_a, tokens)
    out.write_text(tex, encoding="utf-8")

def render_pdflatex(wd, out_path, fonts_dir, repo_root):
    # compile inside the workdir (in-tree paths satisfy tectonic's sandbox)
    data, tokens = _latex_data(wd, repo_root)
    fdir, adir = _stage(fonts_dir, Path(wd) / "assets", Path(wd) / "_tex")
    tex = latexout.to_latex(data, "fonts", "assets" if adir else "", tokens)  # relative to _tex/ (tectonic cwd)
    tex_path = Path(wd) / "_tex" / "_compile.tex"; tex_path.write_text(tex, encoding="utf-8")
    ok, msg = latexout.compile_pdf(tex_path, out_path)
    if not ok:
        raise RuntimeError("pdf-latex compile: " + msg)

RENDERERS = {"pdf": render_pdf, "epub": render_epub, "docx": render_docx, "pptx": render_pptx,
             "md": render_md, "latex": render_latex, "pdf-latex": render_pdflatex}
EXT = {"md": "md", "latex": "tex", "pdf-latex": "pdf"}

if __name__ == "__main__":
    wd, fmt = sys.argv[1], sys.argv[2]; o = opts(sys.argv[3:])
    fonts_dir = o.get("fonts") or str((HERE / "../../runtime/translate/fonts").resolve())
    repo_root = o.get("repo-root")
    meta, _, _ = load(wd)
    code = meta["to"]["code"]
    name = o.get("name") or "document"
    outdir = o.get("out-dir") or str(Path(wd).resolve())
    out_path = str(Path(outdir) / f"{name}.{code}.{EXT.get(fmt, fmt)}")
    RENDERERS[fmt](wd, out_path, fonts_dir, repo_root)
    print(json.dumps({"fmt": fmt, "out": out_path}))
